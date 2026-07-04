"""PPTX extractor: one Markdown section per slide."""
from __future__ import annotations

from pathlib import Path

from .base import ExtractionResult


def extract(src: Path, _assets_dir: Path) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        return ExtractionResult(
            status="manual_review",
            extractor="pptx",
            markdown="",
            error=f"python-pptx missing: {exc}",
        )

    try:
        prs = Presentation(str(src))
    except Exception as exc:  # noqa: BLE001
        return ExtractionResult(
            status="manual_review",
            extractor="pptx",
            markdown="",
            error=f"open failed: {exc}",
        )

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]
    except ImportError:
        MSO_SHAPE_TYPE = None  # older python-pptx; group recursion degrades

    parts: list[str] = []
    counts = {"pictures": 0, "charts": 0}
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide) or f"Slide {i}"
        parts.append(f"## {i}. {title}")
        body_lines: list[str] = []
        _collect_shapes(slide.shapes, body_lines, counts, MSO_SHAPE_TYPE)
        if body_lines:
            parts.append("\n".join(body_lines))
        # Use has_notes_slide, not `slide.notes_slide` truthiness: accessing
        # .notes_slide CREATES a notes part when absent (mutating the deck in
        # memory) and is always truthy, so the guard never short-circuits.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_txt = slide.notes_slide.notes_text_frame.text.strip()
            if notes_txt:
                parts.append(f"**Notes:** {notes_txt}")
        parts.append("")

    # Pictures and charts carry no extractable text, so a deck with them is
    # only partially captured. Be honest and downgrade.
    notes: list[str] = []
    status = "processed"
    dropped = counts["pictures"] + counts["charts"]
    if dropped:
        bits = []
        if counts["pictures"]:
            bits.append(f"{counts['pictures']} picture(s)")
        if counts["charts"]:
            bits.append(f"{counts['charts']} chart(s)")
        notes.append(f"{' and '.join(bits)} not extracted (text captured, visuals not)")
        status = "partial"

    return ExtractionResult(
        status=status,
        extractor="pptx",
        markdown="\n".join(parts),
        notes=notes,
    )


def _collect_shapes(shapes, body_lines: list[str], counts: dict, mso) -> None:
    """Walk a shape collection, appending Markdown lines. Recurses into group
    shapes (whose text would otherwise be dropped) and renders native tables;
    counts pictures/charts so the caller can flag partial extraction."""
    for shape in shapes:
        if mso is not None and shape.shape_type == mso.GROUP:
            _collect_shapes(shape.shapes, body_lines, counts, mso)
            continue
        if shape.has_table:
            body_lines.extend(_table_lines(shape.table))
            continue
        if getattr(shape, "has_chart", False):
            counts["charts"] += 1
            continue
        if mso is not None and shape.shape_type == mso.PICTURE:
            counts["pictures"] += 1
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt:
                    body_lines.append(f"- {txt}")


def _cell(text: str) -> str:
    return (text or "").strip().replace("|", "\\|").replace("\n", " ")


def _table_lines(table) -> list[str]:
    """Render a PowerPoint table to Markdown rows."""
    rows = [[_cell(c.text) for c in row.cells] for row in table.rows]
    if not rows:
        return []
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return out


def _slide_title(slide) -> str | None:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip() or None
    return None
