"""Extractor unit tests over small synthetic fixtures.

The extractors package had no coverage despite holding the riskiest
branching in the pipeline (fence collisions, table escaping, honesty
downgrades). These pin the correctness fixes without needing MinerU,
torch, or a vision API.
"""
from __future__ import annotations

import base64
from pathlib import Path


from ingest_lib.extractors.base import fence
from ingest_lib.extractors import dataset as ds
from ingest_lib.extractors import text as text_ex


# --------------------------------------------------------------- fence()

def test_fence_plain_uses_three_backticks():
    out = fence("hello", "py")
    assert out == "```py\nhello\n```"


def test_fence_grows_past_inner_backtick_run():
    # Content containing ``` must be wrapped in a LONGER fence, else it
    # closes early.
    content = "before\n```\ninner\n```\nafter"
    out = fence(content, "md")
    assert out.startswith("````md\n")
    assert out.endswith("\n````")
    # The inner ``` no longer terminates the block.
    assert out.count("````") == 2


# ---------------------------------------------------------------- text.py

def test_text_extractor_survives_embedded_code_fence(tmp_path: Path):
    src = tmp_path / "note.md"
    src.write_text("# Title\n\n```python\nprint('x')\n```\n", encoding="utf-8")
    res = text_ex.extract(src, tmp_path / "assets")
    assert res.status == "processed"
    # The wrapping fence is longer than the inner ``` so nothing spills out.
    assert res.markdown.startswith("````md\n")
    assert "print('x')" in res.markdown


def test_text_extractor_byte_cap_is_true_bytes(tmp_path: Path, monkeypatch):
    monkeypatch.setattr(text_ex, "_MAX_BYTES", 8)
    src = tmp_path / "big.txt"
    src.write_text("abcdefghijklmnop", encoding="utf-8")  # 16 bytes
    res = text_ex.extract(src, tmp_path / "a")
    assert res.status == "partial"
    assert any("truncated" in n for n in res.notes)


# ------------------------------------------------------------- dataset.py

def test_dataset_csv_escapes_pipe_in_header(tmp_path: Path):
    src = tmp_path / "d.csv"
    src.write_text("price|usd,name\n10,widget\n", encoding="utf-8")
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "processed"
    # The '|' in the header is escaped so it can't add a phantom column.
    assert "price\\|usd" in res.markdown


def test_dataset_jsonl_marks_partial_on_unparseable_lines(tmp_path: Path):
    src = tmp_path / "d.jsonl"
    src.write_text('{"a": 1}\nnot json\n{"a": 2}\n', encoding="utf-8")
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "partial"
    assert any("unparseable" in n for n in res.notes)
    assert "**Records:** 2" in res.markdown


def test_dataset_csv_strips_excel_bom_from_header(tmp_path: Path):
    # Excel's "CSV UTF-8" export prepends a BOM; it must not leak into the
    # first column name of the schema table.
    src = tmp_path / "d.csv"
    src.write_bytes(b"\xef\xbb\xbfname,age\nalice,30\n")
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "processed"
    assert "| 1 | `name` |" in res.markdown
    assert "﻿" not in res.markdown


def test_dataset_jsonl_bom_does_not_drop_first_record(tmp_path: Path):
    src = tmp_path / "d.jsonl"
    src.write_bytes(b'\xef\xbb\xbf{"a": 1}\n{"a": 2}\n')
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "processed"
    assert res.notes == []
    assert "**Records:** 2" in res.markdown


def test_dataset_csv_quoted_crlf_cell_stays_on_one_preview_row(tmp_path: Path):
    # A quoted multiline cell (Excel Alt+Enter) reaches _clip with \r\n
    # intact; a bare CR surviving into the note splits the table row.
    src = tmp_path / "d.csv"
    src.write_bytes(b'name,notes\r\nalice,"line1\r\nline2"\r\nbob,plain\r\n')
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "processed"
    assert "\r" not in res.markdown
    assert "| alice | line1 line2 |" in res.markdown


def test_dataset_csv_oversized_field_is_manual_review(tmp_path: Path):
    import csv
    src = tmp_path / "d.csv"
    big = "x" * (csv.field_size_limit() + 10)
    src.write_text(f'col\n"{big}"\n', encoding="utf-8")
    res = ds.extract(src, tmp_path / "a")
    assert res.status == "manual_review"
    assert "csv parse failed" in (res.error or "")


# ------------------------------------------------------- docx (D8/F110)

# 1x1 transparent PNG for embedding.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def test_docx_extracts_table_and_flags_embedded_image(tmp_path: Path):
    from docx import Document
    from docx.oxml import OxmlElement
    doc = Document()
    doc.add_paragraph("Intro paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "price|usd"   # pipe must be escaped
    table.cell(0, 1).text = "name"
    table.cell(1, 0).text = "10"
    table.cell(1, 1).text = "widget"
    # Inject a w:drawing element (what the extractor counts) rather than a
    # real image, so the test doesn't depend on python-docx's image parser.
    run = doc.add_paragraph("figure below").add_run()
    run._r.append(OxmlElement("w:drawing"))
    src = tmp_path / "d.docx"
    doc.save(str(src))

    from ingest_lib.extractors import docx as docx_ex
    res = docx_ex.extract(src, tmp_path / "a")
    assert "Intro paragraph." in res.markdown
    assert "price\\|usd" in res.markdown            # table extracted + escaped
    assert "| widget |" in res.markdown
    # An embedded image can't be extracted -> honest partial + note.
    assert res.status == "partial"
    assert any("image" in n or "drawing" in n for n in res.notes)


# ------------------------------------------------------- pptx (D8/F036)

def test_pptx_extracts_table_and_flags_picture(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "A bullet of text"
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "col|a"
    tbl.cell(0, 1).text = "colb"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    import io
    slide.shapes.add_picture(io.BytesIO(_PNG), Inches(1), Inches(3))
    src = tmp_path / "p.pptx"
    prs.save(str(src))

    from ingest_lib.extractors import pptx as pptx_ex
    res = pptx_ex.extract(src, tmp_path / "a")
    assert "A bullet of text" in res.markdown
    assert "col\\|a" in res.markdown                # table cell extracted + escaped
    assert "| colb |" in res.markdown
    assert res.status == "partial"                  # picture present
    assert any("picture" in n for n in res.notes)


# ------------------------------------------------------------------ vlm.py

def _gemini_resp(text, finish_reason):
    from types import SimpleNamespace
    return SimpleNamespace(
        text=text, candidates=[SimpleNamespace(finish_reason=finish_reason)]
    )


def _patch_gemini(monkeypatch, resp):
    from google import genai
    monkeypatch.setenv("GOOGLE_API_KEY", "test-key")

    class FakeModels:
        def generate_content(self, **_kw):
            return resp

    class FakeClient:
        def __init__(self, **_kw):
            self.models = FakeModels()

    monkeypatch.setattr(genai, "Client", FakeClient)


def test_vlm_gemini_blank_page_is_success_not_failure(monkeypatch):
    # Contract (stated three times in vlm.py): "" == blank page (success),
    # None only on failure. The gemini helper must not coerce "" to None.
    from google.genai import types
    from ingest_lib.extractors import vlm

    _patch_gemini(monkeypatch, _gemini_resp("", types.FinishReason.STOP))
    res = vlm._vision_gemini(png=b"png", model="gemini-2.5-flash")
    assert res is not None
    assert res.text == ""
    assert res.truncated is False


def test_vlm_gemini_reports_max_tokens_truncation(monkeypatch):
    from google.genai import types
    from ingest_lib.extractors import vlm

    _patch_gemini(
        monkeypatch, _gemini_resp("cut off mid", types.FinishReason.MAX_TOKENS)
    )
    res = vlm._vision_gemini(png=b"png", model="gemini-2.5-flash")
    assert res is not None
    assert res.truncated is True


def test_vlm_anthropic_reports_max_tokens_truncation(monkeypatch):
    import anthropic
    from types import SimpleNamespace
    from ingest_lib.extractors import vlm

    block = anthropic.types.TextBlock.model_construct(type="text", text="cut off")
    resp = SimpleNamespace(content=[block], stop_reason="max_tokens")

    class FakeMessages:
        def create(self, **_kw):
            return resp

    class FakeClient:
        def __init__(self, **_kw):
            self.messages = FakeMessages()

    monkeypatch.setattr(anthropic, "Anthropic", FakeClient)
    res = vlm._vision_anthropic(png=b"png", model="claude-sonnet-4-6")
    assert res is not None
    assert res.text == "cut off"
    assert res.truncated is True


def test_vlm_openai_reports_length_truncation(monkeypatch):
    import openai
    from types import SimpleNamespace
    from ingest_lib.extractors import vlm

    resp = SimpleNamespace(
        choices=[
            SimpleNamespace(
                message=SimpleNamespace(content="cut off"), finish_reason="length"
            )
        ]
    )

    class FakeCompletions:
        def create(self, **_kw):
            return resp

    class FakeClient:
        def __init__(self, **_kw):
            self.chat = SimpleNamespace(completions=FakeCompletions())

    monkeypatch.setattr(openai, "OpenAI", FakeClient)
    res = vlm._vision_openai_compatible(
        png=b"png", model="gpt-5-mini", provider="openai"
    )
    assert res is not None
    assert res.truncated is True


def test_vlm_extract_marks_truncated_pages_partial(monkeypatch, tmp_path: Path):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import vlm

    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    monkeypatch.setattr(vlm, "_render_pages", lambda _src: [b"p1", b"p2"])

    def fake_transcribe(*, png, provider, model, page_no):
        if page_no == 1:
            return vlm._VisionText("dense page", truncated=True)
        return vlm._VisionText("fine page", truncated=False)

    monkeypatch.setattr(vlm, "_transcribe_page", fake_transcribe)
    res = vlm.extract(tmp_path / "doc.pdf", tmp_path / "assets")
    assert res.status == "partial"
    assert any("truncated" in n for n in res.notes)
    # Truncated is NOT failed: the all-pages-failed cleanup must not fire and
    # the truncated text must survive with a visible marker.
    assert not any("failed transcription" in n for n in res.notes)
    assert "dense page" in res.markdown
    assert "truncated at the model output cap" in res.markdown
    assert len(res.assets) == 2


def test_vlm_extract_blank_page_stays_processed(monkeypatch, tmp_path: Path):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import vlm

    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    monkeypatch.setattr(vlm, "_render_pages", lambda _src: [b"p1"])
    monkeypatch.setattr(
        vlm,
        "_transcribe_page",
        lambda **_kw: vlm._VisionText("", truncated=False),
    )
    res = vlm.extract(tmp_path / "doc.pdf", tmp_path / "assets")
    assert res.status == "processed"
    assert "_(blank page)_" in res.markdown


def test_vlm_extract_all_failed_reports_first_cause(monkeypatch, tmp_path: Path):
    # When every page fails, the manual_review error must carry the first
    # page's underlying cause, not just a bare count — otherwise the failure
    # is undiagnosable from index.jsonl / the log.
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import vlm

    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    monkeypatch.setattr(vlm, "_render_pages", lambda _src: [b"p1", b"p2"])
    monkeypatch.setattr(
        vlm,
        "_transcribe_page",
        lambda **_kw: vlm._PageFailure("RuntimeError('vision api exploded')"),
    )
    assets = tmp_path / "assets"
    res = vlm.extract(tmp_path / "doc.pdf", assets)
    assert res.status == "manual_review"
    assert res.error is not None
    assert "all 2 page(s) failed transcription" in res.error
    assert "vision api exploded" in res.error   # the threaded cause
    # Existing behaviour preserved: assets cleaned up, none returned.
    assert not assets.exists()
    assert res.assets == []


def test_vlm_extract_local_model_env_read_at_call_time(monkeypatch, tmp_path: Path):
    # BRAIN_LOCAL_MODEL is read when extract() runs, not at import — setting it
    # after the module is imported must still be honoured (matches summarize).
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import vlm

    monkeypatch.setattr(_summ, "_select_provider", lambda: "local")
    monkeypatch.setattr(vlm, "_render_pages", lambda _src: [b"p1"])
    monkeypatch.delenv("BRAIN_VLM_MODEL", raising=False)
    monkeypatch.setenv("BRAIN_LOCAL_MODEL", "my-local-vision:custom")

    seen: dict[str, str] = {}

    def fake_transcribe(*, png, provider, model, page_no):
        seen["model"] = model
        return vlm._VisionText("ok", truncated=False)

    monkeypatch.setattr(vlm, "_transcribe_page", fake_transcribe)
    vlm.extract(tmp_path / "doc.pdf", tmp_path / "assets")
    assert seen["model"] == "my-local-vision:custom"


def test_mineru_output_handling_crash_degrades_to_manual_review(monkeypatch, tmp_path: Path):
    # An unexpected exception while gathering MinerU's outputs must degrade to
    # manual_review (so extract() falls back to pypdf), not abort ingestion.
    from types import SimpleNamespace
    from ingest_lib.extractors import pdf

    monkeypatch.setattr(
        pdf.subprocess, "run",
        lambda *a, **k: SimpleNamespace(returncode=0, stdout="", stderr=""),
    )

    def boom(*_a, **_k):
        raise RuntimeError("locate exploded")

    monkeypatch.setattr(pdf, "_locate_mineru_outputs", boom)
    res = pdf._extract_with_mineru(tmp_path / "x.pdf", tmp_path / "assets")
    assert res.status == "manual_review"
    assert res.error is not None
    assert "output handling failed" in res.error
    assert "locate exploded" in res.error


# ------------------------------------------------------------------ image.py

def _png_bytes() -> bytes:
    from PIL import Image
    import io
    buf = io.BytesIO()
    Image.new("RGB", (8, 8), (255, 0, 0)).save(buf, format="PNG")
    return buf.getvalue()


def test_image_extension_is_registered():
    from ingest_lib.extractors import dispatch_extractor
    from ingest_lib.extractors import image as image_ex
    from pathlib import Path
    assert dispatch_extractor(Path("photo.jpg")) is image_ex.extract
    assert dispatch_extractor(Path("scan.HEIC")) is image_ex.extract


def test_image_no_vision_provider_is_manual_review(tmp_path: Path, monkeypatch):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import image as image_ex
    monkeypatch.setattr(_summ, "_select_provider", lambda: None)
    src = tmp_path / "x.png"
    src.write_bytes(_png_bytes())
    res = image_ex.extract(src, tmp_path / "a")
    assert res.status == "manual_review"
    assert "no vision" in (res.error or "").lower()


def test_image_transcription_becomes_processed(tmp_path: Path, monkeypatch):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import image as image_ex
    from ingest_lib.extractors import vlm
    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    monkeypatch.setattr(
        vlm, "_transcribe_page",
        lambda **_kw: vlm._VisionText("# Whiteboard\n\nSprint goals", truncated=False),
    )
    src = tmp_path / "board.jpg"
    src.write_bytes(_png_bytes())
    res = image_ex.extract(src, tmp_path / "a")
    assert res.status == "processed"
    assert "Sprint goals" in res.markdown


def test_image_transcription_failure_is_manual_review(tmp_path: Path, monkeypatch):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import image as image_ex
    from ingest_lib.extractors import vlm
    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    monkeypatch.setattr(
        vlm, "_transcribe_page",
        lambda **_kw: vlm._PageFailure("rate limited"),
    )
    src = tmp_path / "x.png"
    src.write_bytes(_png_bytes())
    res = image_ex.extract(src, tmp_path / "a")
    assert res.status == "manual_review"
    assert "rate limited" in (res.error or "")


def test_image_undecodable_file_is_manual_review(tmp_path: Path, monkeypatch):
    from ingest_lib import summarize as _summ
    from ingest_lib.extractors import image as image_ex
    monkeypatch.setattr(_summ, "_select_provider", lambda: "anthropic")
    src = tmp_path / "corrupt.png"
    src.write_bytes(b"this is not a real png")
    res = image_ex.extract(src, tmp_path / "a")
    assert res.status == "manual_review"
    assert "decode" in (res.error or "").lower()
